# """
# PPTX Extractor for HoloLearn
# Extracts text content from PowerPoint presentations.
# """

# from pathlib import Path
# from typing import Dict, Any, Optional, List
# from datetime import datetime
# import json
# import time
# import re
# import tempfile
# import shutil

# import sys
# sys.path.append(str(Path(__file__).parent.parent))
# from utils.configs import OUTPUT_DIR, LOGS_DIR
# from utils.error_handler import ErrorHandler
# from utils.text_cleaner import TextCleaner

# try:
#     from pptx import Presentation
#     _PPTX_AVAILABLE = True
# except ImportError:
#     _PPTX_AVAILABLE = False

# try:
#     import win32com.client
#     _WIN32COM_AVAILABLE = True
# except ImportError:
#     _WIN32COM_AVAILABLE = False

# # Shape type constants (python-pptx / OOXML)
# _SHAPE_PICTURE = 13
# _SHAPE_MEDIA   = 15
# _SHAPE_CHART   = 3
# _SHAPE_TABLE   = 17


# class PPTXExtractor:
#     """Extract text from PowerPoint files"""

#     def __init__(self):
#         if not _PPTX_AVAILABLE:
#             raise ImportError(
#                 "python-pptx not installed. Install with: pip install python-pptx"
#             )

#         self.text_cleaner = TextCleaner()
#         self.base_output_dir = OUTPUT_DIR
#         self.base_logs_dir = LOGS_DIR

#         self.base_output_dir.mkdir(parents=True, exist_ok=True)
#         self.base_logs_dir.mkdir(parents=True, exist_ok=True)

#     def _convert_ppt_to_pptx(self, ppt_path: Path) -> Optional[Path]:
#         """Convert legacy .ppt → .pptx via PowerPoint COM or LibreOffice."""
#         tmp_dir = Path(tempfile.mkdtemp())
#         pptx_path = tmp_dir / (ppt_path.stem + ".pptx")

#         if _WIN32COM_AVAILABLE:
#             try:
#                 import pythoncom
#                 pythoncom.CoInitialize()
#                 ppt_app = win32com.client.Dispatch("PowerPoint.Application")
#                 ppt_app.Visible = 1
#                 presentation = ppt_app.Presentations.Open(
#                     str(ppt_path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False
#                 )
#                 presentation.SaveAs(str(pptx_path.resolve()), 24)  # 24 = ppSaveAsOpenXMLPresentation
#                 presentation.Close()
#                 ppt_app.Quit()
#                 if pptx_path.exists():
#                     return pptx_path
#             except Exception:
#                 pass
#             finally:
#                 try:
#                     pythoncom.CoUninitialize()
#                 except Exception:
#                     pass

#         import subprocess
#         for cmd in ("libreoffice", "soffice"):
#             try:
#                 result = subprocess.run(
#                     [cmd, "--headless", "--convert-to", "pptx",
#                      "--outdir", str(tmp_dir), str(ppt_path.resolve())],
#                     capture_output=True, timeout=60
#                 )
#                 if result.returncode == 0 and pptx_path.exists():
#                     return pptx_path
#             except (FileNotFoundError, subprocess.TimeoutExpired):
#                 continue

#         shutil.rmtree(tmp_dir, ignore_errors=True)
#         return None

#     def _create_resource_name(self, filename: str) -> str:
#         name = Path(filename).stem.lower()
#         name = re.sub(r'[^\w\s-]', '', name)
#         name = re.sub(r'[-\s]+', '_', name).strip('_')
#         return (name[:50] if len(name) > 50 else name) or "unnamed_resource"

#     def _setup_resource_directories(self, resource_name: str, output_dir_override: Optional[Path] = None) -> tuple:
#         resource_output_dir = Path(output_dir_override) if output_dir_override else self.base_output_dir / resource_name
#         resource_logs_dir = self.base_logs_dir / resource_name
#         resource_output_dir.mkdir(parents=True, exist_ok=True)
#         resource_logs_dir.mkdir(parents=True, exist_ok=True)
#         return resource_output_dir, resource_logs_dir

#     # ------------------------------------------------------------------ #
#     #  Structured extraction helpers                                       #
#     # ------------------------------------------------------------------ #

#     def _detect_layout_type(self, slide, slide_num: int, total_slides: int) -> str:
#         """Classify a slide's purpose from its layout name, title text, and content."""
#         layout_name = ""
#         try:
#             layout_name = (slide.slide_layout.name or "").lower()
#         except Exception:
#             pass

#         title_text = self._get_slide_title(slide).lower()

#         summary_kw = ["summary", "conclusion", "recap", "q&a", "questions", "thank", "end"]
#         agenda_kw  = ["agenda", "outline", "contents", "overview", "topics", "table of"]
#         title_kw   = ["title", "cover", "opening", "intro"]

#         if slide_num == 1 or any(k in layout_name for k in title_kw):
#             return "title_slide"
#         if any(k in title_text for k in summary_kw):
#             return "summary_slide"
#         if any(k in title_text for k in agenda_kw):
#             return "agenda_slide"
#         if slide_num == total_slides and any(k in title_text for k in summary_kw):
#             return "summary_slide"

#         pic_count = sum(1 for s in slide.shapes if getattr(s, 'shape_type', None) == _SHAPE_PICTURE)
#         text_shapes = sum(1 for s in slide.shapes if hasattr(s, 'text') and s.text.strip())
#         if pic_count > 0 and text_shapes <= 1:
#             return "image_slide"

#         return "content_slide"

#     def _get_slide_title(self, slide) -> str:
#         """Return the title placeholder text, or first text shape text, or empty string."""
#         for shape in slide.shapes:
#             if not shape.has_text_frame:
#                 continue
#             ph = getattr(shape, 'placeholder_format', None)
#             if ph is not None and ph.idx == 0:  # idx 0 = TITLE placeholder
#                 return shape.text_frame.text.strip()
#         return ""

#     def _sorted_shapes(self, slide) -> list:
#         """Return slide shapes sorted by vertical then horizontal position (reading order)."""
#         def _key(s):
#             return (s.top if s.top is not None else 0,
#                     s.left if s.left is not None else 0)
#         return sorted(slide.shapes, key=_key)

#     def _extract_slide_data(self, slide, slide_num: int, total_slides: int,
#                             include_notes: bool) -> dict:
#         """
#         Return a structured dict for one slide:
#         {slide_number, title, body (list of {text, level}), notes, layout_type, has_media}
#         """
#         title = self._get_slide_title(slide)
#         layout_type = self._detect_layout_type(slide, slide_num, total_slides)

#         has_media = any(
#             getattr(s, 'shape_type', None) in (_SHAPE_PICTURE, _SHAPE_MEDIA, _SHAPE_CHART)
#             for s in slide.shapes
#         )

#         body_items: List[dict] = []
#         for shape in self._sorted_shapes(slide):
#             if not shape.has_text_frame:
#                 continue
#             ph = getattr(shape, 'placeholder_format', None)
#             if ph is not None and ph.idx == 0:
#                 continue  # already captured as title
#             for para in shape.text_frame.paragraphs:
#                 text = para.text.strip()
#                 if text:
#                     body_items.append({"text": text, "level": para.level})

#         notes_text = ""
#         if include_notes and slide.has_notes_slide:
#             try:
#                 notes_text = slide.notes_slide.notes_text_frame.text.strip()
#             except Exception:
#                 pass

#         return {
#             "slide_number": slide_num,
#             "title": title,
#             "body": body_items,
#             "notes": notes_text,
#             "layout_type": layout_type,
#             "has_media": has_media,
#         }

#     def _compute_quality_score(self, text: str, sections: List[dict]) -> float:
#         word_score = min(len(text.split()) / 5000, 1.0)
#         structure_score = min(len(sections) / 10, 1.0)
#         return round(word_score * 0.6 + structure_score * 0.4, 2)

#     # ------------------------------------------------------------------ #
#     #  Main extraction                                                     #
#     # ------------------------------------------------------------------ #

#     def extract(self,
#                 pptx_path: str,
#                 resource_id: Optional[str] = None,
#                 clean_text: bool = True,
#                 include_notes: bool = True,
#                 output_dir: Optional[str] = None) -> Dict[str, Any]:
#         """
#         Extract text from a PowerPoint file.

#         Returns dict with: success, resource_name, resource_id, text_file,
#         metadata_file, structured_file, output_dir, logs_dir, extracted_text,
#         slides_data, sections, content_quality_score, metadata.
#         """
#         start_time = time.time()
#         pptx_path = Path(pptx_path)
#         _tmp_dir_to_cleanup = None

#         if pptx_path.suffix.lower() == ".ppt":
#             converted = self._convert_ppt_to_pptx(pptx_path)
#             if converted is None:
#                 resource_name = self._create_resource_name(pptx_path.name)
#                 override = Path(output_dir) if output_dir else None
#                 out_dir, _ = self._setup_resource_directories(resource_name, output_dir_override=override)
#                 return self._create_error_result(
#                     resource_name,
#                     "Cannot convert .ppt file: install Microsoft Office (win32com) or LibreOffice.",
#                     out_dir, pptx_path.name
#                 )
#             _tmp_dir_to_cleanup = converted.parent
#             pptx_path = converted

#         resource_name = self._create_resource_name(pptx_path.name)
#         override = Path(output_dir) if output_dir else None
#         output_dir, logs_dir = self._setup_resource_directories(resource_name, output_dir_override=override)

#         error_handler = ErrorHandler(f"pptx_{resource_name}")
#         error_handler.log_file = logs_dir / "extraction.log"
#         error_handler.logger = error_handler._setup_logger()

#         if not pptx_path.exists():
#             error_msg = f"PPTX file not found: {pptx_path}"
#             error_handler.log_error(FileNotFoundError(error_msg), context="Validating PPTX file",
#                                     metadata={"path": str(pptx_path)})
#             return self._create_error_result(resource_name, error_msg, output_dir, pptx_path.name)

#         file_size = pptx_path.stat().st_size
#         file_size_mb = file_size / (1024 * 1024)

#         error_handler.log_info(f"Starting PPTX extraction: {pptx_path.name}",
#                                metadata={"size_mb": f"{file_size_mb:.2f}",
#                                          "resource_name": resource_name,
#                                          "output_dir": str(output_dir)})

#         try:
#             prs = Presentation(pptx_path)
#             total_slides = len(prs.slides)
#             slides_data: List[dict] = []
#             slides_with_notes = 0
#             media_slide_count = 0

#             for slide_num, slide in enumerate(prs.slides, 1):
#                 slide_dict = self._extract_slide_data(slide, slide_num, total_slides, include_notes)
#                 slides_data.append(slide_dict)
#                 if slide_dict["notes"]:
#                     slides_with_notes += 1
#                 if slide_dict["has_media"]:
#                     media_slide_count += 1

#             # Build flat extracted_text from slides_data (reading-order guaranteed)
#             extracted_text = ""
#             for sd in slides_data:
#                 extracted_text += f"\n{'='*60}\n"
#                 extracted_text += f"SLIDE {sd['slide_number']} [{sd['layout_type']}]\n"
#                 extracted_text += f"{'='*60}\n\n"

#                 if sd["title"]:
#                     extracted_text += sd["title"] + "\n\n"

#                 for item in sd["body"]:
#                     indent = "  " * item["level"]
#                     extracted_text += f"{indent}{item['text']}\n"

#                 if not sd["title"] and not sd["body"]:
#                     extracted_text += "[No text content on this slide]\n"

#                 if sd["notes"]:
#                     extracted_text += f"\n[SPEAKER NOTES]\n{sd['notes']}\n"

#                 extracted_text += "\n"

#             if clean_text:
#                 extracted_text = self.text_cleaner.clean_text(
#                     extracted_text, remove_urls=False, remove_emails=False, fix_spacing=True
#                 )
#             extracted_text = self.text_cleaner.remove_duplicate_lines(extracted_text)

#             # Derive lecture-ready sections from slides (one section per slide with a title)
#             sections = [
#                 {
#                     "title": sd["title"] or f"Slide {sd['slide_number']}",
#                     "body": " ".join(item["text"] for item in sd["body"]),
#                     "type": sd["layout_type"],
#                     "source_location": {"slide": sd["slide_number"]},
#                 }
#                 for sd in slides_data
#             ]

#             quality_score = self._compute_quality_score(extracted_text, sections)
#             processing_time = time.time() - start_time

#             metadata = {
#                 "resource_name": resource_name,
#                 "resource_id": resource_id or resource_name,
#                 "filename": pptx_path.name,
#                 "source_type": "pptx",
#                 "upload_date": datetime.now().isoformat(),
#                 "extraction_timestamp": datetime.now().isoformat(),
#                 "file_size_bytes": file_size,
#                 "processing_time_seconds": round(processing_time, 2),
#                 "status": "success",
#                 "error_message": None,
#                 "slide_count": total_slides,
#                 "slides_with_notes": slides_with_notes,
#                 "media_slide_count": media_slide_count,
#                 "character_count": len(extracted_text),
#                 "included_notes": include_notes,
#                 "content_quality_score": quality_score,
#             }

#             text_file = output_dir / "text.txt"
#             text_file.write_text(extracted_text, encoding='utf-8')
#             metadata["extracted_text_path"] = str(text_file)

#             metadata_file = output_dir / "metadata.json"
#             metadata_file.write_text(json.dumps(metadata, indent=2), encoding='utf-8')

#             structured = {"slides": slides_data, "sections": sections}
#             structured_file = output_dir / "structured.json"
#             structured_file.write_text(json.dumps(structured, indent=2, ensure_ascii=False), encoding='utf-8')

#             error_handler.log_success(f"PPTX extracted successfully: {pptx_path.name}",
#                                       metadata={"slides": total_slides, "media_slides": media_slide_count,
#                                                 "time": f"{processing_time:.2f}s"})

#             return {
#                 "success": True,
#                 "resource_name": resource_name,
#                 "resource_id": resource_id or resource_name,
#                 "text_file": str(text_file),
#                 "metadata_file": str(metadata_file),
#                 "structured_file": str(structured_file),
#                 "output_dir": str(output_dir),
#                 "logs_dir": str(logs_dir),
#                 "extracted_text": extracted_text,
#                 "slides_data": slides_data,
#                 "sections": sections,
#                 "content_quality_score": quality_score,
#                 "metadata": metadata,
#             }

#         except Exception as e:
#             processing_time = time.time() - start_time
#             error_handler.log_error(e, context=f"Extracting PPTX: {pptx_path.name}",
#                                     metadata={"resource_name": resource_name})
#             return self._create_error_result(resource_name, str(e), output_dir,
#                                              pptx_path.name, file_size, processing_time)
#         finally:
#             if _tmp_dir_to_cleanup:
#                 shutil.rmtree(_tmp_dir_to_cleanup, ignore_errors=True)

#     def _create_error_result(self,
#                              resource_name: str,
#                              error_message: str,
#                              output_dir: Path,
#                              filename: str = "unknown",
#                              file_size: int = 0,
#                              processing_time: float = 0) -> Dict[str, Any]:
#         metadata = {
#             "resource_name": resource_name,
#             "filename": filename,
#             "source_type": "pptx",
#             "upload_date": datetime.now().isoformat(),
#             "extraction_timestamp": datetime.now().isoformat(),
#             "file_size_bytes": file_size,
#             "processing_time_seconds": round(processing_time, 2),
#             "status": "failed",
#             "error_message": error_message,
#         }
#         metadata_file = output_dir / "metadata.json"
#         metadata_file.write_text(json.dumps(metadata, indent=2), encoding='utf-8')
#         return {
#             "success": False,
#             "resource_name": resource_name,
#             "text_file": None,
#             "metadata_file": str(metadata_file),
#             "structured_file": None,
#             "output_dir": str(output_dir),
#             "extracted_text": "",
#             "slides_data": [],
#             "sections": [],
#             "content_quality_score": 0.0,
#             "metadata": metadata,
#             "error": error_message,
#         }

#     def extract_metadata_only(self, pptx_path: str) -> Dict[str, Any]:
#         """Extract only metadata without extracting text."""
#         _tmp_dir = None
#         try:
#             pptx_path = Path(pptx_path)
#             if pptx_path.suffix.lower() == ".ppt":
#                 converted = self._convert_ppt_to_pptx(pptx_path)
#                 if converted is None:
#                     return {}
#                 _tmp_dir = converted.parent
#                 pptx_path = converted

#             prs = Presentation(pptx_path)
#             slides_with_notes = sum(1 for s in prs.slides if s.has_notes_slide)
#             return {
#                 "filename": pptx_path.name,
#                 "resource_name": self._create_resource_name(pptx_path.name),
#                 "slide_count": len(prs.slides),
#                 "slides_with_notes": slides_with_notes,
#                 "file_size_bytes": pptx_path.stat().st_size,
#                 "file_size_mb": round(pptx_path.stat().st_size / (1024 * 1024), 2),
#                 "core_properties": {
#                     "title": prs.core_properties.title or "N/A",
#                     "author": prs.core_properties.author or "N/A",
#                     "subject": prs.core_properties.subject or "N/A",
#                     "created": str(prs.core_properties.created) if prs.core_properties.created else "N/A",
#                     "modified": str(prs.core_properties.modified) if prs.core_properties.modified else "N/A",
#                 },
#             }
#         except Exception as e:
#             ErrorHandler("pptx_metadata").log_error(e, context=f"Extracting metadata from {pptx_path}")
#             return {}
#         finally:
#             if _tmp_dir:
#                 shutil.rmtree(_tmp_dir, ignore_errors=True)


# # Example usage and testing
# if __name__ == "__main__":
#     from utils.file_picker import FilePicker

#     print("=== Testing PPTX Extractor ===\n")

#     extractor = PPTXExtractor()

#     picker = FilePicker()
#     print("Please select a PowerPoint file...")
#     test_pptx = picker.pick_pptx()
#     picker.close()

#     if test_pptx:
#         print(f"\n✓ Selected: {Path(test_pptx).name}\n")

#         print("1. Extracting metadata only...")
#         metadata = extractor.extract_metadata_only(test_pptx)
#         print(f"   Slides: {metadata.get('slide_count', 'N/A')}")
#         print(f"   Slides with notes: {metadata.get('slides_with_notes', 'N/A')}\n")

#         print("2. Full extraction (with speaker notes)...")
#         result = extractor.extract(pptx_path=test_pptx, clean_text=True, include_notes=True)

#         if result['success']:
#             print(f"   ✓ Success!")
#             print(f"   Slides: {result['metadata']['slide_count']}")
#             print(f"   Media slides: {result['metadata']['media_slide_count']}")
#             print(f"   Quality score: {result['content_quality_score']}")
#             print(f"   Structured file: {result['structured_file']}")

#             print("\n   Slide layout breakdown:")
#             layout_counts: Dict[str, int] = {}
#             for sd in result['slides_data']:
#                 lt = sd['layout_type']
#                 layout_counts[lt] = layout_counts.get(lt, 0) + 1
#             for lt, count in layout_counts.items():
#                 print(f"     {lt}: {count}")
#         else:
#             print(f"   ✗ Failed: {result['error']}")
#     else:
#         print("❌ No file selected")
#         print("   The extractor is ready to use!")


"""
PPTX Extractor for HoloLearn
Extracts text content from PowerPoint presentations.

Modes
-----
extract()        — python-pptx only (text slides, fast, exact)
extract_hybrid() — per-slide decision: python-pptx for text slides,
                   CCA pipeline for image-heavy slides.
                   Output format matches extract() exactly.
"""

from pathlib import Path
from typing import Dict, Any, Optional, List
from datetime import datetime
import json
import time
import re
import tempfile
import shutil
import subprocess

import sys
sys.path.append(str(Path(__file__).parent.parent))
from utils.configs import (
    OUTPUT_DIR, LOGS_DIR,
    OCR_PIPELINE_DIR, CNN_MODEL_PATH, CLASSIFIER_DEVICE,
    VLM_MODEL, VLM_BASE_URL, VLM_API_KEY,
)
from utils.error_handler import ErrorHandler
from utils.text_cleaner import TextCleaner

try:
    from pptx import Presentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

try:
    import win32com.client
    _WIN32COM_AVAILABLE = True
except ImportError:
    _WIN32COM_AVAILABLE = False


# Shape type constants (python-pptx / OOXML)
_SHAPE_PICTURE = 13
_SHAPE_MEDIA   = 15
_SHAPE_CHART   = 3
_SHAPE_TABLE   = 17


class PPTXExtractor:
    """Extract text from PowerPoint files."""

    def __init__(self):
        if not _PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )

        self.text_cleaner    = TextCleaner()
        self.base_output_dir = OUTPUT_DIR
        self.base_logs_dir   = LOGS_DIR

        self.base_output_dir.mkdir(parents=True, exist_ok=True)
        self.base_logs_dir.mkdir(parents=True, exist_ok=True)

        # Inject the PARENT of OCR_PIPELINE_DIR so OCR is importable as a package.
        # (OCR has __init__.py — injecting the dir itself breaks relative imports)
        if OCR_PIPELINE_DIR:
            _ocr_parent = str(Path(OCR_PIPELINE_DIR).parent)
            if _ocr_parent not in sys.path:
                sys.path.insert(0, _ocr_parent)

    # ------------------------------------------------------------------ #
    #  Legacy .ppt conversion                                              #
    # ------------------------------------------------------------------ #

    def _convert_ppt_to_pptx(self, ppt_path: Path) -> Optional[Path]:
        """Convert legacy .ppt → .pptx via PowerPoint COM or LibreOffice."""
        tmp_dir   = Path(tempfile.mkdtemp())
        pptx_path = tmp_dir / (ppt_path.stem + ".pptx")

        if _WIN32COM_AVAILABLE:
            try:
                import pythoncom
                pythoncom.CoInitialize()
                ppt_app      = win32com.client.Dispatch("PowerPoint.Application")
                ppt_app.Visible = 1
                presentation = ppt_app.Presentations.Open(
                    str(ppt_path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False
                )
                presentation.SaveAs(str(pptx_path.resolve()), 24)
                presentation.Close()
                ppt_app.Quit()
                if pptx_path.exists():
                    return pptx_path
            except Exception:
                pass
            finally:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

        for cmd in ("libreoffice", "soffice"):
            try:
                result = subprocess.run(
                    [cmd, "--headless", "--convert-to", "pptx",
                     "--outdir", str(tmp_dir), str(ppt_path.resolve())],
                    capture_output=True, timeout=60,
                )
                if result.returncode == 0 and pptx_path.exists():
                    return pptx_path
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue

        shutil.rmtree(tmp_dir, ignore_errors=True)
        return None

    # ------------------------------------------------------------------ #
    #  Directory / naming helpers                                          #
    # ------------------------------------------------------------------ #

    def _create_resource_name(self, filename: str) -> str:
        name = Path(filename).stem.lower()
        name = re.sub(r'[^\w\s-]', '', name)
        name = re.sub(r'[-\s]+', '_', name).strip('_')
        return (name[:50] if len(name) > 50 else name) or "unnamed_resource"

    def _setup_resource_directories(
        self,
        resource_name: str,
        output_dir_override: Optional[Path] = None,
    ) -> tuple:
        resource_output_dir = (
            Path(output_dir_override) if output_dir_override
            else self.base_output_dir / resource_name
        )
        resource_logs_dir = self.base_logs_dir / resource_name
        resource_output_dir.mkdir(parents=True, exist_ok=True)
        resource_logs_dir.mkdir(parents=True, exist_ok=True)
        return resource_output_dir, resource_logs_dir

    # ------------------------------------------------------------------ #
    #  Structured extraction helpers                                       #
    # ------------------------------------------------------------------ #

    def _detect_layout_type(self, slide, slide_num: int, total_slides: int) -> str:
        layout_name = ""
        try:
            layout_name = (slide.slide_layout.name or "").lower()
        except Exception:
            pass

        title_text = self._get_slide_title(slide).lower()
        summary_kw = ["summary", "conclusion", "recap", "q&a", "questions", "thank", "end"]
        agenda_kw  = ["agenda", "outline", "contents", "overview", "topics", "table of"]
        title_kw   = ["title", "cover", "opening", "intro"]

        if slide_num == 1 or any(k in layout_name for k in title_kw):
            return "title_slide"
        if any(k in title_text for k in summary_kw):
            return "summary_slide"
        if any(k in title_text for k in agenda_kw):
            return "agenda_slide"
        if slide_num == total_slides and any(k in title_text for k in summary_kw):
            return "summary_slide"

        pic_count   = sum(1 for s in slide.shapes if getattr(s, 'shape_type', None) == _SHAPE_PICTURE)
        text_shapes = sum(1 for s in slide.shapes if hasattr(s, 'text') and s.text.strip())
        if pic_count > 0 and text_shapes <= 1:
            return "image_slide"

        return "content_slide"

    def _get_slide_title(self, slide) -> str:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:
                    return shape.text_frame.text.strip()
            except Exception:
                pass
        return ""

    def _sorted_shapes(self, slide) -> list:
        def _key(s):
            return (s.top if s.top is not None else 0,
                    s.left if s.left is not None else 0)
        return sorted(slide.shapes, key=_key)

    def _extract_slide_data(
        self, slide, slide_num: int, total_slides: int, include_notes: bool
    ) -> dict:
        title       = self._get_slide_title(slide)
        layout_type = self._detect_layout_type(slide, slide_num, total_slides)

        has_media = any(
            getattr(s, 'shape_type', None) in (_SHAPE_PICTURE, _SHAPE_MEDIA, _SHAPE_CHART)
            for s in slide.shapes
        )

        body_items: List[dict] = []
        for shape in self._sorted_shapes(slide):
            if not shape.has_text_frame:
                continue
            try:
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:
                    continue  # already captured as title
            except Exception:
                pass
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    body_items.append({"text": text, "level": para.level})

        notes_text = ""
        if include_notes and slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

        return {
            "slide_number": slide_num,
            "title":        title,
            "body":         body_items,
            "notes":        notes_text,
            "layout_type":  layout_type,
            "has_media":    has_media,
        }

    def _compute_quality_score(self, text: str, sections: List[dict]) -> float:
        word_score      = min(len(text.split()) / 5000, 1.0)
        structure_score = min(len(sections) / 10, 1.0)
        return round(word_score * 0.6 + structure_score * 0.4, 2)

    # ------------------------------------------------------------------ #
    #  Hybrid extraction helpers                                           #
    # ------------------------------------------------------------------ #

    def _slide_has_text(self, slide_dict: dict, min_chars: int = 30) -> bool:
        """True if the slide has enough embedded text to skip CCA."""
        all_text = slide_dict["title"] + " ".join(
            item["text"] for item in slide_dict["body"]
        )
        return len(all_text.strip()) >= min_chars

    # ------------------------------------------------------------------ #
    #  extract()  — python-pptx only                                      #
    # ------------------------------------------------------------------ #

    def extract(
        self,
        pptx_path:     str,
        resource_id:   Optional[str] = None,
        clean_text:    bool = True,
        include_notes: bool = True,
        output_dir:    Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Extract text from a PowerPoint file using python-pptx only.
        Best for presentations with text-based slides.
        """
        start_time          = time.time()
        pptx_path           = Path(pptx_path)
        _tmp_dir_to_cleanup = None

        if pptx_path.suffix.lower() == ".ppt":
            converted = self._convert_ppt_to_pptx(pptx_path)
            if converted is None:
                resource_name = self._create_resource_name(pptx_path.name)
                override      = Path(output_dir) if output_dir else None
                out_dir, _    = self._setup_resource_directories(
                    resource_name, output_dir_override=override
                )
                return self._create_error_result(
                    resource_name,
                    "Cannot convert .ppt: install Microsoft Office or LibreOffice.",
                    out_dir, pptx_path.name,
                )
            _tmp_dir_to_cleanup = converted.parent
            pptx_path           = converted

        resource_name = self._create_resource_name(pptx_path.name)
        override      = Path(output_dir) if output_dir else None
        output_dir, logs_dir = self._setup_resource_directories(
            resource_name, output_dir_override=override
        )

        error_handler          = ErrorHandler(f"pptx_{resource_name}")
        error_handler.log_file = logs_dir / "extraction.log"
        error_handler.logger   = error_handler._setup_logger()

        if not pptx_path.exists():
            return self._create_error_result(
                resource_name, f"PPTX not found: {pptx_path}", output_dir, pptx_path.name
            )

        file_size    = pptx_path.stat().st_size
        file_size_mb = file_size / (1024 * 1024)

        error_handler.log_info(
            f"Starting PPTX extraction: {pptx_path.name}",
            metadata={"size_mb": f"{file_size_mb:.2f}", "resource_name": resource_name},
        )

        try:
            prs               = Presentation(pptx_path)
            total_slides      = len(prs.slides)
            slides_data: List[dict] = []
            slides_with_notes = 0
            media_slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                sd = self._extract_slide_data(slide, slide_num, total_slides, include_notes)
                slides_data.append(sd)
                if sd["notes"]:
                    slides_with_notes += 1
                if sd["has_media"]:
                    media_slide_count += 1

            extracted_text = ""
            for sd in slides_data:
                extracted_text += f"\n{'='*60}\n"
                extracted_text += f"SLIDE {sd['slide_number']} [{sd['layout_type']}]\n"
                extracted_text += f"{'='*60}\n\n"
                if sd["title"]:
                    extracted_text += sd["title"] + "\n\n"
                for item in sd["body"]:
                    extracted_text += "  " * item["level"] + item["text"] + "\n"
                if not sd["title"] and not sd["body"]:
                    extracted_text += "[No text content on this slide]\n"
                if sd["notes"]:
                    extracted_text += f"\n[SPEAKER NOTES]\n{sd['notes']}\n"
                extracted_text += "\n"

            if clean_text:
                extracted_text = self.text_cleaner.clean_text(
                    extracted_text, remove_urls=False, remove_emails=False, fix_spacing=True
                )
            extracted_text = self.text_cleaner.remove_duplicate_lines(extracted_text)

            sections = [
                {
                    "title":           sd["title"] or f"Slide {sd['slide_number']}",
                    "body":            " ".join(item["text"] for item in sd["body"]),
                    "type":            sd["layout_type"],
                    "source_location": {"slide": sd["slide_number"]},
                }
                for sd in slides_data
            ]

            quality_score   = self._compute_quality_score(extracted_text, sections)
            processing_time = time.time() - start_time

            return self._save_outputs(
                resource_name     = resource_name,
                resource_id       = resource_id or resource_name,
                pptx_path         = pptx_path,
                output_dir        = output_dir,
                logs_dir          = logs_dir,
                extracted_text    = extracted_text,
                slides_data       = slides_data,
                sections          = sections,
                quality_score     = quality_score,
                file_size         = file_size,
                processing_time   = processing_time,
                total_slides      = total_slides,
                slides_with_notes = slides_with_notes,
                media_slide_count = media_slide_count,
                include_notes     = include_notes,
                error_handler     = error_handler,
                extraction_mode   = "pptx",
            )

        except Exception as e:
            processing_time = time.time() - start_time
            error_handler.log_error(e, context=f"Extracting PPTX: {pptx_path.name}")
            return self._create_error_result(
                resource_name, str(e), output_dir,
                pptx_path.name, file_size, processing_time,
            )
        finally:
            if _tmp_dir_to_cleanup:
                shutil.rmtree(_tmp_dir_to_cleanup, ignore_errors=True)

    # ------------------------------------------------------------------ #
    #  extract_hybrid()  — python-pptx for text, CCA for image slides     #
    # ------------------------------------------------------------------ #

    def _extract_slide_images(
        self, slide, output_dir: Path, slide_num: int
    ) -> List[str]:
        """
        Extract all embedded picture shapes from a slide and save to disk.

        Uses short sequential filenames (s0001_p00.png) to stay well under
        the Windows MAX_PATH limit.

        Returns list of saved image paths (only successfully saved ones).
        """
        output_dir.mkdir(parents=True, exist_ok=True)
        img_paths = []

        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type != _SHAPE_PICTURE:
                continue
            try:
                image    = shape.image
                ext      = (image.ext or "png").lstrip(".")
                img_path = output_dir / f"s{slide_num:04d}_p{idx:03d}.{ext}"
                img_path.write_bytes(image.blob)
                img_paths.append(str(img_path))
            except Exception:
                pass   # skip linked / unreadable images silently

        return img_paths

    # ------------------------------------------------------------------ #
    #  extract_hybrid()  — python-pptx for text, CCA for image slides     #
    # ------------------------------------------------------------------ #

    def extract_hybrid(
        self,
        pptx_path:         str,
        resource_id:       Optional[str] = None,
        output_dir:        Optional[str] = None,
        include_notes:     bool          = True,
        clean_text:        bool          = True,
        min_text_chars:    int           = 30,
        # ── OCR pipeline params — default to configs.py values ───────────
        model_path:        Optional[str] = None,
        vlm_model:         Optional[str] = None,
        vlm_base_url:      Optional[str] = None,
        vlm_api_key:       Optional[str] = None,
        classifier_device: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Hybrid extraction — per-slide decision:
          text slide  (>= min_text_chars embedded text) → python-pptx (fast, exact)
          image slide (mostly visual, little/no text)   → extract pictures → CCA pipeline

        No external dependencies needed — uses python-pptx's built-in image access.
        Output format is IDENTICAL to extract() so callers need no changes.
        Extra metadata keys: text_slides, image_slides, extraction_mode="hybrid".
        """
        # ── Resolve OCR params: call-level override → config default ──────
        _model_path   = model_path        or CNN_MODEL_PATH
        _vlm_model    = vlm_model         or VLM_MODEL
        _vlm_base_url = vlm_base_url      or VLM_BASE_URL
        _api_key      = vlm_api_key       or VLM_API_KEY
        _device       = classifier_device or CLASSIFIER_DEVICE

        start_time          = time.time()
        pptx_path           = Path(pptx_path)
        _tmp_dir_to_cleanup = None

        # ── Handle legacy .ppt ────────────────────────────────────────────
        if pptx_path.suffix.lower() == ".ppt":
            converted = self._convert_ppt_to_pptx(pptx_path)
            if converted is None:
                resource_name = self._create_resource_name(pptx_path.name)
                override      = Path(output_dir) if output_dir else None
                out_dir, _    = self._setup_resource_directories(
                    resource_name, output_dir_override=override
                )
                return self._create_error_result(
                    resource_name,
                    "Cannot convert .ppt: install Microsoft Office or LibreOffice.",
                    out_dir, pptx_path.name,
                )
            _tmp_dir_to_cleanup = converted.parent
            pptx_path           = converted

        resource_name = self._create_resource_name(pptx_path.name)
        override      = Path(output_dir) if output_dir else None
        out_dir, logs_dir = self._setup_resource_directories(
            resource_name, output_dir_override=override
        )

        error_handler          = ErrorHandler(f"pptx_hybrid_{resource_name}")
        error_handler.log_file = logs_dir / "extraction.log"
        error_handler.logger   = error_handler._setup_logger()

        if not pptx_path.exists():
            return self._create_error_result(
                resource_name, f"PPTX not found: {pptx_path}", out_dir, pptx_path.name
            )

        file_size    = pptx_path.stat().st_size
        file_size_mb = file_size / (1024 * 1024)

        # ── Import OCR pipeline ───────────────────────────────────────────
        try:
            from OCR.pipeline import run_full_pipeline
        except ImportError as e:
            error_handler.log_error(
                e, context="Importing OCR pipeline — is OCR_PIPELINE_DIR set correctly?"
            )
            return self._create_error_result(
                resource_name, f"OCR pipeline import failed: {e}",
                out_dir, pptx_path.name, file_size,
            )

        try:
            # ── Step 1: Parse all slides with python-pptx ─────────────────
            prs               = Presentation(pptx_path)
            total_slides      = len(prs.slides)
            slides_data: List[dict] = []
            slides_with_notes = 0
            media_slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                sd = self._extract_slide_data(slide, slide_num, total_slides, include_notes)
                slides_data.append(sd)
                if sd["notes"]:
                    slides_with_notes += 1
                if sd["has_media"]:
                    media_slide_count += 1

            # ── Step 2: Per-slide decision ─────────────────────────────────
            all_text_parts = []
            text_slides    = []
            image_slides   = []

            for sd, slide in zip(slides_data, prs.slides):
                slide_num   = sd["slide_number"]
                slide_label = f"SLIDE {slide_num} [{sd['layout_type']}]"

                if self._slide_has_text(sd, min_text_chars):
                    # ── Text slide: use python-pptx output ────────────────
                    text_slides.append(slide_num)

                    part = f"\n{'='*60}\n{slide_label}\n{'='*60}\n\n"
                    if sd["title"]:
                        part += sd["title"] + "\n\n"
                    for item in sd["body"]:
                        part += "  " * item["level"] + item["text"] + "\n"
                    if not sd["title"] and not sd["body"]:
                        part += "[No text content on this slide]\n"
                    if sd["notes"]:
                        part += f"\n[SPEAKER NOTES]\n{sd['notes']}\n"
                    all_text_parts.append(part)

                else:
                    # ── Image slide: extract pictures → CCA ───────────────
                    image_slides.append(slide_num)
                    error_handler.log_info(f"Slide {slide_num} is image-heavy → CCA pipeline")

                    imgs_dir  = out_dir / f"s{slide_num:04d}_imgs"
                    img_paths = self._extract_slide_images(slide, imgs_dir, slide_num)

                    slide_texts = []

                    for img_path in img_paths:
                        # Short crop dir name to avoid MAX_PATH
                        crops_dir = str(out_dir / f"s{slide_num:04d}_{Path(img_path).stem}_crops")
                        try:
                            docs, _, _ = run_full_pipeline(
                                image_path        = img_path,
                                model_path        = _model_path,
                                api_key           = _api_key,
                                output_dir        = crops_dir,
                                classifier_device = _device,
                                vlm_model         = _vlm_model,
                                vlm_base_url      = _vlm_base_url,
                            )
                            if docs and docs[0]["text"].strip():
                                slide_texts.append(docs[0]["text"])
                        except Exception as e:
                            error_handler.log_error(
                                e, context=f"CCA on slide {slide_num} "
                                           f"image {Path(img_path).name}"
                            )

                    # Cleanup extracted picture files
                    for img_path in img_paths:
                        try:
                            Path(img_path).unlink(missing_ok=True)
                        except OSError:
                            pass
                    try:
                        imgs_dir.rmdir()   # only removes if empty
                    except OSError:
                        pass

                    if not slide_texts:
                        # No pictures found or all CCA failed — python-pptx fallback
                        fallback = " ".join(item["text"] for item in sd["body"])
                        slide_texts = [fallback or "[Image slide — no extractable content]"]

                    cca_text = "\n\n".join(slide_texts)

                    part = f"\n{'='*60}\n{slide_label}\n{'='*60}\n\n"
                    if sd["title"]:
                        part += sd["title"] + "\n\n"
                    part += cca_text + "\n"
                    if sd["notes"]:
                        part += f"\n[SPEAKER NOTES]\n{sd['notes']}\n"
                    all_text_parts.append(part)

            # ── Step 3: Assemble and clean text ────────────────────────────
            extracted_text = "\n".join(all_text_parts)

            if clean_text:
                extracted_text = self.text_cleaner.clean_text(
                    extracted_text, remove_urls=False, remove_emails=False, fix_spacing=True
                )
            extracted_text = self.text_cleaner.remove_duplicate_lines(extracted_text)

            sections = [
                {
                    "title":           sd["title"] or f"Slide {sd['slide_number']}",
                    "body":            " ".join(item["text"] for item in sd["body"]),
                    "type":            sd["layout_type"],
                    "source_location": {"slide": sd["slide_number"]},
                }
                for sd in slides_data
            ]

            quality_score   = self._compute_quality_score(extracted_text, sections)
            processing_time = time.time() - start_time

            error_handler.log_info(
                f"Hybrid complete: {len(text_slides)} text slides, "
                f"{len(image_slides)} image slides"
            )

            return self._save_outputs(
                resource_name     = resource_name,
                resource_id       = resource_id or resource_name,
                pptx_path         = pptx_path,
                output_dir        = out_dir,
                logs_dir          = logs_dir,
                extracted_text    = extracted_text,
                slides_data       = slides_data,
                sections          = sections,
                quality_score     = quality_score,
                file_size         = file_size,
                processing_time   = processing_time,
                total_slides      = total_slides,
                slides_with_notes = slides_with_notes,
                media_slide_count = media_slide_count,
                include_notes     = include_notes,
                error_handler     = error_handler,
                extraction_mode   = "hybrid",
                text_slides       = text_slides,
                image_slides      = image_slides,
            )

        except Exception as e:
            processing_time = time.time() - start_time
            error_handler.log_error(e, context=f"Hybrid extracting PPTX: {pptx_path.name}")
            return self._create_error_result(
                resource_name, str(e), out_dir,
                pptx_path.name, file_size, processing_time,
            )
        finally:
            if _tmp_dir_to_cleanup:
                shutil.rmtree(_tmp_dir_to_cleanup, ignore_errors=True)

    def _save_outputs(
        self,
        resource_name:     str,
        resource_id:       str,
        pptx_path:         Path,
        output_dir:        Path,
        logs_dir:          Path,
        extracted_text:    str,
        slides_data:       List[dict],
        sections:          List[dict],
        quality_score:     float,
        file_size:         int,
        processing_time:   float,
        total_slides:      int,
        slides_with_notes: int,
        media_slide_count: int,
        include_notes:     bool,
        error_handler,
        extraction_mode:   str                  = "pptx",
        text_slides:       Optional[List[int]]  = None,
        image_slides:      Optional[List[int]]  = None,
    ) -> Dict[str, Any]:
        """
        Write text / metadata / structured JSON and return the standard result dict.
        Shared by both extract() and extract_hybrid() — output format is identical.
        """
        metadata = {
            "resource_name":           resource_name,
            "resource_id":             resource_id,
            "filename":                pptx_path.name,
            "source_type":             "pptx",
            "extraction_mode":         extraction_mode,
            "upload_date":             datetime.now().isoformat(),
            "extraction_timestamp":    datetime.now().isoformat(),
            "file_size_bytes":         file_size,
            "processing_time_seconds": round(processing_time, 2),
            "status":                  "success",
            "error_message":           None,
            "slide_count":             total_slides,
            "slides_with_notes":       slides_with_notes,
            "media_slide_count":       media_slide_count,
            "character_count":         len(extracted_text),
            "included_notes":          include_notes,
            "content_quality_score":   quality_score,
        }
        if text_slides is not None:
            metadata["text_slides"]  = text_slides
        if image_slides is not None:
            metadata["image_slides"] = image_slides

        text_file = output_dir / f"{resource_name}_text.txt"
        text_file.write_text(extracted_text, encoding="utf-8")
        metadata["extracted_text_path"] = str(text_file)

        metadata_file = output_dir / f"{resource_name}_metadata.json"
        metadata_file.write_text(json.dumps(metadata, indent=2), encoding="utf-8")

        structured_file = output_dir / f"{resource_name}_structured.json"
        structured_file.write_text(
            json.dumps({"slides": slides_data, "sections": sections},
                       indent=2, ensure_ascii=False),
            encoding="utf-8",
        )

        error_handler.log_success(
            f"PPTX extracted: {pptx_path.name}",
            metadata={
                "slides": total_slides,
                "mode":   extraction_mode,
                "time":   f"{processing_time:.2f}s",
            },
        )

        return {
            "success":               True,
            "resource_name":         resource_name,
            "resource_id":           resource_id,
            "text_file":             str(text_file),
            "metadata_file":         str(metadata_file),
            "structured_file":       str(structured_file),
            "output_dir":            str(output_dir),
            "logs_dir":              str(logs_dir),
            "extracted_text":        extracted_text,
            "slides_data":           slides_data,
            "sections":              sections,
            "content_quality_score": quality_score,
            "metadata":              metadata,
            "text_slides":           text_slides,
            "image_slides":          image_slides,
        }

    # ------------------------------------------------------------------ #
    #  Error result                                                        #
    # ------------------------------------------------------------------ #

    def _create_error_result(
        self,
        resource_name:   str,
        error_message:   str,
        output_dir:      Path,
        filename:        str   = "unknown",
        file_size:       int   = 0,
        processing_time: float = 0,
    ) -> Dict[str, Any]:
        metadata = {
            "resource_name":           resource_name,
            "filename":                filename,
            "source_type":             "pptx",
            "upload_date":             datetime.now().isoformat(),
            "extraction_timestamp":    datetime.now().isoformat(),
            "file_size_bytes":         file_size,
            "processing_time_seconds": round(processing_time, 2),
            "status":                  "failed",
            "error_message":           error_message,
        }
        metadata_file = output_dir / "metadata.json"
        metadata_file.write_text(json.dumps(metadata, indent=2), encoding="utf-8")
        return {
            "success":               False,
            "resource_name":         resource_name,
            "text_file":             None,
            "metadata_file":         str(metadata_file),
            "structured_file":       None,
            "output_dir":            str(output_dir),
            "logs_dir":              None,
            "extracted_text":        "",
            "slides_data":           [],
            "sections":              [],
            "content_quality_score": 0.0,
            "metadata":              metadata,
            "text_slides":           None,
            "image_slides":          None,
            "error":                 error_message,
        }

    # ------------------------------------------------------------------ #
    #  Metadata only                                                       #
    # ------------------------------------------------------------------ #

    def extract_metadata_only(self, pptx_path: str) -> Dict[str, Any]:
        """Extract only metadata without extracting text."""
        _tmp_dir = None
        try:
            pptx_path = Path(pptx_path)
            if pptx_path.suffix.lower() == ".ppt":
                converted = self._convert_ppt_to_pptx(pptx_path)
                if converted is None:
                    return {}
                _tmp_dir  = converted.parent
                pptx_path = converted

            prs               = Presentation(pptx_path)
            slides_with_notes = sum(1 for s in prs.slides if s.has_notes_slide)
            return {
                "filename":          pptx_path.name,
                "resource_name":     self._create_resource_name(pptx_path.name),
                "slide_count":       len(prs.slides),
                "slides_with_notes": slides_with_notes,
                "file_size_bytes":   pptx_path.stat().st_size,
                "file_size_mb":      round(pptx_path.stat().st_size / (1024 * 1024), 2),
                "core_properties": {
                    "title":    prs.core_properties.title    or "N/A",
                    "author":   prs.core_properties.author   or "N/A",
                    "subject":  prs.core_properties.subject  or "N/A",
                    "created":  str(prs.core_properties.created)  if prs.core_properties.created  else "N/A",
                    "modified": str(prs.core_properties.modified) if prs.core_properties.modified else "N/A",
                },
            }
        except Exception as e:
            ErrorHandler("pptx_metadata").log_error(
                e, context=f"Extracting metadata from {pptx_path}"
            )
            return {}
        finally:
            if _tmp_dir:
                shutil.rmtree(_tmp_dir, ignore_errors=True)


# ================================================================== #
#  Quick test                                                          #
# ================================================================== #

if __name__ == "__main__":
    from utils.file_picker import FilePicker

    print("=== Testing PPTX Extractor ===\n")
    extractor = PPTXExtractor()

    picker    = FilePicker()
    print("Select a PowerPoint file...")
    test_pptx = picker.pick_pptx()
    picker.close()

    if not test_pptx:
        print("No file selected.")
    else:
        print(f"\n✓ Selected: {Path(test_pptx).name}\n")

        print("1. Metadata only...")
        meta = extractor.extract_metadata_only(test_pptx)
        print(f"   Slides: {meta.get('slide_count')}  |  Notes slides: {meta.get('slides_with_notes')}")

        print("\n2. python-pptx extraction...")
        r = extractor.extract(pptx_path=test_pptx, clean_text=True, include_notes=True)
        if r["success"]:
            print(f"   ✓  Text file:     {r['text_file']}")
            print(f"      Slides:        {r['metadata']['slide_count']}")
            print(f"      Media slides:  {r['metadata']['media_slide_count']}")
            print(f"      Quality score: {r['content_quality_score']}")
        else:
            print(f"   ✗  {r['error']}")

        print("\n3. Hybrid extraction (CCA for image-heavy slides)...")
        rh = extractor.extract_hybrid(pptx_path=test_pptx, clean_text=True, include_notes=True)
        if rh["success"]:
            print(f"   ✓  Text file:     {rh['text_file']}")
            print(f"      Text slides:   {rh['text_slides']}")
            print(f"      Image slides:  {rh['image_slides']}")
            print(f"      Quality score: {rh['content_quality_score']}")
        else:
            print(f"   ✗  {rh['error']}")